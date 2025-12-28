"""
Docling document parser for advanced document processing.

Docling provides high-quality document parsing with:
- Structure-aware PDF parsing
- Table extraction with structure preservation
- Image extraction and description
- Multi-format support (PDF, DOCX, PPTX, HTML, etc.)

Usage:
    from app.parsing.parsers.docling_parser import DoclingParser

    parser = DoclingParser()
    documents = parser.parse(Path("document.pdf"))
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from langchain_core.documents import Document

from app.core.config import settings

logger = logging.getLogger(__name__)


# Configuration
DOCLING_ENABLED = getattr(settings, "DOCLING_ENABLED", False)
DOCLING_OCR_ENABLED = getattr(settings, "DOCLING_OCR_ENABLED", True)
DOCLING_TABLE_MODE = getattr(settings, "DOCLING_TABLE_MODE", "markdown")  # markdown | html | plain


class DoclingParser:
    """
    Docling-based document parser.

    Provides high-quality document parsing with structure preservation,
    table extraction, and multi-format support.
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".html", ".md", ".asciidoc"}

    def __init__(
        self,
        ocr_enabled: bool = True,
        table_mode: str = "markdown",
        extract_images: bool = False,
        max_pages: Optional[int] = None,
    ):
        """
        Initialize the Docling parser.

        Args:
            ocr_enabled: Enable OCR for scanned documents
            table_mode: Table output format (markdown, html, plain)
            extract_images: Extract and describe images
            max_pages: Maximum pages to process (None = all)
        """
        self.ocr_enabled = ocr_enabled
        self.table_mode = table_mode
        self.extract_images = extract_images
        self.max_pages = max_pages
        self._converter = None

    def _get_converter(self) -> Any:
        """Lazy initialization of Docling converter."""
        if self._converter is not None:
            return self._converter

        try:
            from docling.document_converter import DocumentConverter
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.datamodel.base_models import InputFormat

            # Configure pipeline options
            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = self.ocr_enabled
            pipeline_options.do_table_structure = True

            self._converter = DocumentConverter(
                allowed_formats=[
                    InputFormat.PDF,
                    InputFormat.DOCX,
                    InputFormat.PPTX,
                    InputFormat.HTML,
                    InputFormat.MD,
                    InputFormat.ASCIIDOC,
                ],
            )
            logger.info("Docling converter initialized")
            return self._converter

        except ImportError:
            logger.error(
                "Docling is not installed. Install with: pip install docling"
            )
            raise ImportError(
                "Docling is not installed. Install with: pip install docling"
            )

    def parse(
        self,
        file_path: Path,
        **kwargs,
    ) -> List[Document]:
        """
        Parse a document using Docling.

        Args:
            file_path: Path to the document
            **kwargs: Additional arguments

        Returns:
            List of parsed Document objects
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if file_path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {file_path.suffix}. "
                f"Supported: {self.SUPPORTED_EXTENSIONS}"
            )

        try:
            return self._parse_with_docling(file_path)
        except ImportError:
            logger.warning("Docling not available, falling back to basic parsing")
            return self._fallback_parse(file_path)

    def _parse_with_docling(self, file_path: Path) -> List[Document]:
        """Parse using Docling library."""
        converter = self._get_converter()

        # Convert document
        result = converter.convert(str(file_path))

        documents = []
        doc = result.document

        # Extract text content
        full_text = doc.export_to_markdown()

        # Create main document
        metadata = {
            "source": str(file_path),
            "filename": file_path.name,
            "parser": "docling",
            "page_count": getattr(doc, "page_count", 1),
        }

        documents.append(Document(
            page_content=full_text,
            metadata=metadata,
        ))

        # Extract tables as separate documents if present
        if hasattr(doc, "tables") and doc.tables:
            for i, table in enumerate(doc.tables):
                table_content = self._format_table(table)
                if table_content:
                    documents.append(Document(
                        page_content=table_content,
                        metadata={
                            **metadata,
                            "content_type": "table",
                            "table_index": i,
                        },
                    ))

        # Extract images if enabled
        if self.extract_images and hasattr(doc, "pictures"):
            for i, pic in enumerate(doc.pictures):
                if hasattr(pic, "caption") and pic.caption:
                    documents.append(Document(
                        page_content=f"[Image {i + 1}]: {pic.caption}",
                        metadata={
                            **metadata,
                            "content_type": "image",
                            "image_index": i,
                        },
                    ))

        logger.info(
            "Docling parsed %s: %d documents extracted",
            file_path.name,
            len(documents),
        )

        return documents

    def _format_table(self, table: Any) -> str:
        """Format a table based on the configured mode."""
        try:
            if self.table_mode == "markdown":
                if hasattr(table, "export_to_markdown"):
                    return table.export_to_markdown()
                return self._table_to_markdown(table)
            elif self.table_mode == "html":
                if hasattr(table, "export_to_html"):
                    return table.export_to_html()
                return self._table_to_html(table)
            else:
                return self._table_to_plain(table)
        except Exception as e:
            logger.warning("Failed to format table: %s", e)
            return ""

    def _table_to_markdown(self, table: Any) -> str:
        """Convert table to markdown format."""
        if not hasattr(table, "data") or not table.data:
            return ""

        rows = table.data
        if not rows:
            return ""

        # Build markdown table
        lines = []

        # Header
        header = rows[0]
        lines.append("| " + " | ".join(str(cell) for cell in header) + " |")
        lines.append("| " + " | ".join("---" for _ in header) + " |")

        # Data rows
        for row in rows[1:]:
            lines.append("| " + " | ".join(str(cell) for cell in row) + " |")

        return "\n".join(lines)

    def _table_to_html(self, table: Any) -> str:
        """Convert table to HTML format."""
        if not hasattr(table, "data") or not table.data:
            return ""

        rows = table.data
        if not rows:
            return ""

        lines = ["<table>"]

        # Header
        lines.append("  <thead><tr>")
        for cell in rows[0]:
            lines.append(f"    <th>{cell}</th>")
        lines.append("  </tr></thead>")

        # Body
        lines.append("  <tbody>")
        for row in rows[1:]:
            lines.append("    <tr>")
            for cell in row:
                lines.append(f"      <td>{cell}</td>")
            lines.append("    </tr>")
        lines.append("  </tbody>")
        lines.append("</table>")

        return "\n".join(lines)

    def _table_to_plain(self, table: Any) -> str:
        """Convert table to plain text format."""
        if not hasattr(table, "data") or not table.data:
            return ""

        rows = table.data
        lines = []

        for row in rows:
            lines.append("\t".join(str(cell) for cell in row))

        return "\n".join(lines)

    def _fallback_parse(self, file_path: Path) -> List[Document]:
        """Fallback parsing when Docling is not available."""
        ext = file_path.suffix.lower()

        if ext == ".pdf":
            return self._fallback_pdf(file_path)
        elif ext in {".docx", ".pptx"}:
            return self._fallback_office(file_path)
        elif ext in {".html", ".md"}:
            return self._fallback_text(file_path)
        else:
            return self._fallback_text(file_path)

    def _fallback_pdf(self, file_path: Path) -> List[Document]:
        """Fallback PDF parsing using PyMuPDF."""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(file_path))
            documents = []

            for page_num in range(len(doc)):
                if self.max_pages and page_num >= self.max_pages:
                    break

                page = doc[page_num]
                text = page.get_text()

                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={
                            "source": str(file_path),
                            "filename": file_path.name,
                            "page": page_num + 1,
                            "parser": "pymupdf_fallback",
                        },
                    ))

            doc.close()
            return documents

        except ImportError:
            logger.error("PyMuPDF not available for fallback")
            return []

    def _fallback_office(self, file_path: Path) -> List[Document]:
        """Fallback Office document parsing."""
        ext = file_path.suffix.lower()

        if ext == ".docx":
            try:
                from docx import Document as DocxDocument

                doc = DocxDocument(str(file_path))
                text = "\n".join(para.text for para in doc.paragraphs)

                return [Document(
                    page_content=text,
                    metadata={
                        "source": str(file_path),
                        "filename": file_path.name,
                        "parser": "python-docx_fallback",
                    },
                )]
            except ImportError:
                logger.error("python-docx not available for fallback")
                return []

        elif ext == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(str(file_path))
                text_parts = []

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)

                return [Document(
                    page_content="\n".join(text_parts),
                    metadata={
                        "source": str(file_path),
                        "filename": file_path.name,
                        "parser": "python-pptx_fallback",
                    },
                )]
            except ImportError:
                logger.error("python-pptx not available for fallback")
                return []

        return []

    def _fallback_text(self, file_path: Path) -> List[Document]:
        """Fallback text file parsing."""
        try:
            text = file_path.read_text(encoding="utf-8")
            return [Document(
                page_content=text,
                metadata={
                    "source": str(file_path),
                    "filename": file_path.name,
                    "parser": "text_fallback",
                },
            )]
        except Exception as e:
            logger.error("Failed to read file: %s", e)
            return []

    async def aparse(
        self,
        file_path: Path,
        **kwargs,
    ) -> List[Document]:
        """Async version of parse."""
        import asyncio
        return await asyncio.to_thread(self.parse, file_path, **kwargs)
