#!/usr/bin/env python3
# ruff: noqa: E402, I001
"""Verify knowledge-base breadth across multiple file types against a live API."""


import argparse
import json
import shutil
import sys
import time
from pathlib import Path
from typing import Any

from pptx import Presentation


def ensure_repo_root_on_sys_path(script_path: str | Path) -> str:
    repo_root = str(Path(script_path).resolve().parents[1])
    if repo_root not in sys.path:
        sys.path.insert(0, repo_root)
    return repo_root


ensure_repo_root_on_sys_path(__file__)

from scripts.remote_kb_boundary_matrix import (
    LiveApi,
    citation_document_ids,
    ensure_success,
    parsed_text_from_response,
    record_step,
    response_text_from_body,
    wait_for_document_completed,
)


DEFAULT_TENANT_ID = "00000000-0000-0000-0000-000000000000"
REPO_ROOT = Path(__file__).resolve().parents[1]
WORD_FIXTURE = REPO_ROOT / "tests/fixtures/parsing_golden_broader/word_project_brief_docx/input/sample.docx"
XLSX_FIXTURE = REPO_ROOT / "tests/fixtures/parsing_golden_broader/excel_budget_sheet_xlsx/input/sample.xlsx"
TEXT_FAMILY_CASE_NAMES = {"text_note_txt", "ini_config", "sql_query"}


def evaluate_format_case(
    case: dict[str, Any],
    *,
    document_id: str,
    chunk_count: int,
    parsed_chars: int,
    citation_doc_ids: list[str],
    citation_count: int,
    response_text: str,
) -> list[str]:
    failures: list[str] = []
    name = str(case.get("name") or "case")
    expected_document_id = str(case.get("expected_document_id") or "").strip()
    expected_terms = [str(item) for item in (case.get("expected_terms") or []) if str(item).strip()]
    min_chunks = int(case.get("min_chunks") or 0)
    min_parsed_chars = int(case.get("min_parsed_chars") or 0)
    min_citations = int(case.get("min_citations") or 0)

    if min_chunks and int(chunk_count) < min_chunks:
        failures.append(f"{name}: min_chunks={min_chunks} actual={int(chunk_count)}")
    if min_parsed_chars and int(parsed_chars) < min_parsed_chars:
        failures.append(f"{name}: min_parsed_chars={min_parsed_chars} actual={int(parsed_chars)}")
    if min_citations and int(citation_count) < min_citations:
        failures.append(f"{name}: min_citations={min_citations} actual={int(citation_count)}")
    if expected_document_id and expected_document_id not in citation_doc_ids:
        failures.append(f"{name}: expected_document_id={expected_document_id} actual={citation_doc_ids}")

    lowered = str(response_text or "").casefold()
    missing_terms = [term for term in expected_terms if term.casefold() not in lowered]
    if missing_terms:
        failures.append(f"{name}: expected_terms missing={missing_terms}")
    return failures


def prepare_fixture_files(fixtures_dir: Path) -> list[dict[str, Any]]:
    fixtures_dir.mkdir(parents=True, exist_ok=True)

    txt_path = fixtures_dir / "kb-note.txt"
    txt_path.write_text(
        "Token TXT-BEACON belongs only to this text note.\n"
        "Owner: Talia Beacon.\n",
        encoding="utf-8",
    )

    ini_path = fixtures_dir / "kb-service.ini"
    ini_path.write_text(
        "[service]\n"
        "token=INI-HARBOR\n"
        "owner=Inez Harbor\n"
        "status=ready\n",
        encoding="utf-8",
    )

    sql_path = fixtures_dir / "kb-owner.sql"
    sql_path.write_text(
        "-- owner lookup\n"
        'SELECT "SQL-LANTERN" AS token, "Soren Lantern" AS owner;\n',
        encoding="utf-8",
    )

    toml_path = fixtures_dir / "kb-service.toml"
    toml_path.write_text(
        'token = "TOML-EMBER"\n'
        'owner = "Toma Ember"\n'
        'status = "armed"\n',
        encoding="utf-8",
    )

    properties_path = fixtures_dir / "kb-service.properties"
    properties_path.write_text(
        "token=PROP-LATTICE\n"
        "owner=Priya Lattice\n"
        "status=warm\n",
        encoding="utf-8",
    )

    env_path = fixtures_dir / "kb-service.env"
    env_path.write_text(
        "TOKEN=ENV-SIGNAL\n"
        "OWNER=Evan Signal\n"
        "STATUS=green\n",
        encoding="utf-8",
    )

    rst_path = fixtures_dir / "kb-guide.rst"
    rst_path.write_text(
        "RST Guide\n"
        "=========\n\n"
        "Token RST-FLARE belongs only to this rst note.\n\n"
        "Owner: Rhea Flare.\n",
        encoding="utf-8",
    )

    log_path = fixtures_dir / "kb-service.log"
    log_path.write_text(
        "2026-05-24T08:00:00Z token=LOG-BEACON owner=Logan Beacon status=hot\n",
        encoding="utf-8",
    )

    jsonl_path = fixtures_dir / "kb-events.jsonl"
    jsonl_path.write_text(
        '{"token":"JSONL-COMET","owner":"Jill Comet","status":"queued"}\n',
        encoding="utf-8",
    )

    proto_path = fixtures_dir / "kb-service.proto"
    proto_path.write_text(
        'syntax = "proto3";\n'
        "// token PROTO-RIDGE owner Priya Ridge\n"
        "message Probe {}\n",
        encoding="utf-8",
    )

    graphql_path = fixtures_dir / "kb-query.graphql"
    graphql_path.write_text(
        "# token GQL-ORBIT owner Gina Orbit\n"
        "query Probe { viewer { id } }\n",
        encoding="utf-8",
    )

    tf_path = fixtures_dir / "kb-main.tf"
    tf_path.write_text(
        'locals { token = "TF-LANTERN" owner = "Tariq Lantern" }\n',
        encoding="utf-8",
    )

    patch_path = fixtures_dir / "kb-change.patch"
    patch_path.write_text(
        "+++ token PATCH-NOVA owner Nia Nova\n"
        "@@\n"
        "+status=applied\n",
        encoding="utf-8",
    )

    pptx_path = fixtures_dir / "kb-briefing.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "PPTX-ANCHOR"
    slide.placeholders[1].text = "Owner: Paula Anchor\nToken PPTX-ANCHOR belongs only to this deck."
    prs.save(pptx_path)

    ndjson_path = fixtures_dir / "kb-events.ndjson"
    ndjson_path.write_text(
        '{"token":"NDJSON-PEAK","owner":"Nadia Peak","status":"ready"}\n',
        encoding="utf-8",
    )

    hcl_path = fixtures_dir / "kb-service.hcl"
    hcl_path.write_text(
        'token = "HCL-VALE"\n'
        'owner = "Hector Vale"\n'
        'status = "steady"\n',
        encoding="utf-8",
    )

    adoc_path = fixtures_dir / "kb-guide.adoc"
    adoc_path.write_text(
        "= ADOC Guide\n\n"
        "Token ADOC-EMBER belongs only to this adoc note.\n\n"
        "Owner: Ada Ember.\n",
        encoding="utf-8",
    )

    diff_path = fixtures_dir / "kb-change.diff"
    diff_path.write_text(
        "--- a\n"
        "+++ b\n"
        "+token DIFF-SHIFT owner Dario Shift\n",
        encoding="utf-8",
    )

    atom_path = fixtures_dir / "kb-feed.atom"
    atom_path.write_text(
        '<?xml version="1.0" encoding="utf-8"?><feed><entry><title>ATOM-NOVA</title><author><name>Anya Nova</name></author></entry></feed>',
        encoding="utf-8",
    )

    markdown_path = fixtures_dir / "kb-note.md"
    markdown_path.write_text(
        "# KB Markdown Note\n\n"
        "Token MD-ORBIT belongs only to the markdown note.\n\n"
        "Owner: Maya Orbit.\n",
        encoding="utf-8",
    )

    html_path = fixtures_dir / "kb-page.html"
    html_path.write_text(
        "<!doctype html><html><body><h1>KB HTML Page</h1><p>Token HTML-CASCADE belongs only to the HTML page.</p><p>Owner: Hugo Cascade.</p></body></html>",
        encoding="utf-8",
    )

    csv_path = fixtures_dir / "kb-metrics.csv"
    csv_path.write_text(
        "token,owner,status\n"
        "CSV-RIDGE,Rita Ridge,healthy\n",
        encoding="utf-8",
    )

    json_path = fixtures_dir / "kb-faq.json"
    json_path.write_text(
        json.dumps(
            {
                "token": "JSON-QUASAR",
                "owner": "Juno Quasar",
                "answer": "JSON-QUASAR belongs only to the JSON fixture.",
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    yaml_path = fixtures_dir / "kb-config.yaml"
    yaml_path.write_text(
        "token: YAML-CINDER\n"
        "owner: Yara Cinder\n"
        "status: approved\n",
        encoding="utf-8",
    )

    xml_path = fixtures_dir / "kb-catalog.xml"
    xml_path.write_text(
        "<?xml version=\"1.0\" encoding=\"UTF-8\"?>\n"
        "<catalog>\n"
        "  <entry token=\"XML-DELTA\">\n"
        "    <owner>Xenia Delta</owner>\n"
        "    <status>verified</status>\n"
        "  </entry>\n"
        "</catalog>\n",
        encoding="utf-8",
    )

    word_target = fixtures_dir / "word-project-brief.docx"
    xlsx_target = fixtures_dir / "excel-budget-sheet.xlsx"
    shutil.copy2(WORD_FIXTURE, word_target)
    shutil.copy2(XLSX_FIXTURE, xlsx_target)

    return [
        {
            "name": "text_note_txt",
            "path": txt_path,
            "query": "Which token belongs only to this text note?",
            "expected_terms": ["TXT-BEACON", "Talia Beacon"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "ini_config",
            "path": ini_path,
            "query": "Who owns token INI-HARBOR?",
            "expected_terms": ["INI-HARBOR", "Inez Harbor"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "sql_query",
            "path": sql_path,
            "query": "Who owns token SQL-LANTERN?",
            "expected_terms": ["SQL-LANTERN", "Soren Lantern"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "toml_config",
            "path": toml_path,
            "query": "Who owns token TOML-EMBER?",
            "expected_terms": ["TOML-EMBER", "Toma Ember"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "properties_file",
            "path": properties_path,
            "query": "Who owns token PROP-LATTICE?",
            "expected_terms": ["PROP-LATTICE", "Priya Lattice"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "env_file",
            "path": env_path,
            "query": "Who owns token ENV-SIGNAL?",
            "expected_terms": ["ENV-SIGNAL", "Evan Signal"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "rst_note",
            "path": rst_path,
            "query": "Which token belongs only to this rst note?",
            "expected_terms": ["RST-FLARE", "Rhea Flare"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "log_file",
            "path": log_path,
            "query": "Who owns token LOG-BEACON?",
            "expected_terms": ["LOG-BEACON", "Logan Beacon"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "jsonl_feed",
            "path": jsonl_path,
            "query": "Who owns token JSONL-COMET?",
            "expected_terms": ["JSONL-COMET", "Jill Comet"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "proto_schema",
            "path": proto_path,
            "query": "Who owns token PROTO-RIDGE?",
            "expected_terms": ["PROTO-RIDGE", "Priya Ridge"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "graphql_doc",
            "path": graphql_path,
            "query": "Who owns token GQL-ORBIT?",
            "expected_terms": ["GQL-ORBIT", "Gina Orbit"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "terraform_file",
            "path": tf_path,
            "query": "Who owns token TF-LANTERN?",
            "expected_terms": ["TF-LANTERN", "Tariq Lantern"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "patch_file",
            "path": patch_path,
            "query": "Who owns token PATCH-NOVA?",
            "expected_terms": ["PATCH-NOVA", "Nia Nova"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "pptx_deck",
            "path": pptx_path,
            "query": "Which token belongs only to this deck?",
            "expected_terms": ["PPTX-ANCHOR", "Paula Anchor"],
            "parser_backend": "auto",
            "family_group": "office_like_family",
        },
        {
            "name": "ndjson_feed",
            "path": ndjson_path,
            "query": "Who owns token NDJSON-PEAK?",
            "expected_terms": ["NDJSON-PEAK", "Nadia Peak"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "hcl_config",
            "path": hcl_path,
            "query": "Who owns token HCL-VALE?",
            "expected_terms": ["HCL-VALE", "Hector Vale"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "asciidoc_note",
            "path": adoc_path,
            "query": "Which token belongs only to this adoc note?",
            "expected_terms": ["ADOC-EMBER", "Ada Ember"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "diff_file",
            "path": diff_path,
            "query": "Who owns token DIFF-SHIFT?",
            "expected_terms": ["DIFF-SHIFT", "Dario Shift"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "atom_feed",
            "path": atom_path,
            "query": "Who owns token ATOM-NOVA in the Atom feed entry?",
            "expected_terms": ["ATOM-NOVA", "Anya Nova"],
            "parser_backend": "basic",
            "family_group": "text_family",
        },
        {
            "name": "markdown_note",
            "path": markdown_path,
            "query": "Which token belongs only to the markdown note?",
            "expected_terms": ["MD-ORBIT"],
        },
        {
            "name": "html_page",
            "path": html_path,
            "query": "Which token belongs only to the HTML page?",
            "expected_terms": ["HTML-CASCADE"],
        },
        {
            "name": "csv_metrics",
            "path": csv_path,
            "query": "Who owns token CSV-RIDGE?",
            "expected_terms": ["CSV-RIDGE"],
        },
        {
            "name": "json_faq",
            "path": json_path,
            "query": "Which token belongs only to the JSON fixture?",
            "expected_terms": ["JSON-QUASAR"],
        },
        {
            "name": "yaml_config",
            "path": yaml_path,
            "query": "Who owns token YAML-CINDER?",
            "expected_terms": ["YAML-CINDER", "Yara Cinder"],
        },
        {
            "name": "xml_catalog",
            "path": xml_path,
            "query": "Who owns token XML-DELTA in the XML catalog entry?",
            "expected_terms": ["XML-DELTA", "Xenia Delta"],
        },
        {
            "name": "word_project_brief_docx",
            "path": word_target,
            "query": "In the Word project brief, who owns the rollout?",
            "expected_terms": ["Lina Chen"],
        },
        {
            "name": "excel_budget_sheet_xlsx",
            "path": xlsx_target,
            "query": "In the Excel budget sheet, what status belongs to APAC?",
            "expected_terms": ["Review"],
        },
    ]


def select_fixture_cases(
    fixture_cases: list[dict[str, Any]],
    *,
    case_names: list[str],
    include_text_families: bool,
) -> list[dict[str, Any]]:
    requested_names = [str(name).strip() for name in case_names if str(name).strip()]
    if requested_names:
        requested_set = set(requested_names)
        return [case for case in fixture_cases if str(case.get("name") or "") in requested_set]
    if include_text_families:
        return fixture_cases
    return [case for case in fixture_cases if str(case.get("family_group") or "") != "text_family"]


def cleanup_dataset(api: LiveApi, *, steps: list[dict[str, Any]], dataset_id: str) -> dict[str, Any]:
    summary: dict[str, Any] = {"dataset_id": dataset_id}
    resp = api.json("POST", f"/api/v1/datasets/{dataset_id}/purge?dry_run=false&max_delete=1000", payload={})
    record_step(steps, f"cleanup:purge:{dataset_id}", resp)
    if 200 <= resp.status < 300:
        summary["purge_deleted"] = int((resp.body or {}).get("deleted") or 0) if isinstance(resp.body, dict) else 0
    resp = api.json("DELETE", f"/api/v1/datasets/{dataset_id}")
    record_step(steps, f"cleanup:delete_dataset:{dataset_id}", resp)
    summary["delete_dataset_status"] = int(resp.status)
    return summary


def main() -> int:
    parser = argparse.ArgumentParser(description="Run KB file-type breadth verification against a live API.")
    parser.add_argument("--base-url", default="http://127.0.0.1:8000")
    parser.add_argument("--tenant-id", default=DEFAULT_TENANT_ID)
    parser.add_argument("--account-id", default="demo")
    parser.add_argument("--user-id", default="demo")
    parser.add_argument("--artifact-dir", default="")
    parser.add_argument("--timeout", type=int, default=180)
    parser.add_argument("--poll-timeout", type=int, default=300)
    parser.add_argument("--include-text-families", action="store_true")
    parser.add_argument("--case-name", action="append", default=[])
    args = parser.parse_args()

    run_id = time.strftime("%Y%m%d-%H%M%S")
    artifact_dir = Path(args.artifact_dir or f"artifacts/kb-format-matrix/{run_id}").resolve()
    artifact_dir.mkdir(parents=True, exist_ok=True)
    fixture_cases = select_fixture_cases(
        prepare_fixture_files(artifact_dir / "fixtures"),
        case_names=list(args.case_name or []),
        include_text_families=bool(args.include_text_families),
    )
    api = LiveApi(args.base_url, args.tenant_id, args.account_id, args.user_id, args.timeout)

    steps: list[dict[str, Any]] = []
    summary: dict[str, Any] = {
        "ok": False,
        "artifact_dir": str(artifact_dir),
        "base_url": args.base_url,
        "dataset_id": "",
        "cases": [],
    }

    dataset_id = ""

    try:
        resp = api.json("GET", "/api/v1/health")
        record_step(steps, "health", resp)
        ensure_success("health", resp)

        resp = api.json(
            "POST",
            "/api/v1/datasets/",
            payload={
                "name": f"KB Format Matrix {run_id}",
                "description": "Knowledge-base file-type breadth verification.",
                "permission": "all_team_members",
                "default_parser_backend": "auto",
                "default_chunk_strategy": "langchain_recursive",
                "pipeline": {
                    "governance_enabled": True,
                    "persist_parsed_content": True,
                    "persist_parsed_content_max_chars": 200000,
                    "chunk_size": 1200,
                    "chunk_overlap": 120,
                    "chunk_vector_enabled": True,
                    "bm25_index_enabled": True,
                    "kg_enabled": False,
                    "event_vector_enabled": False,
                    "entity_vector_enabled": False,
                },
            },
        )
        record_step(steps, "create_dataset", resp)
        ensure_success("create_dataset", resp)
        dataset_id = str((resp.body or {}).get("id") or (resp.body or {}).get("dataset_id") or "")
        if not dataset_id:
            raise RuntimeError("create_dataset missing dataset id")
        summary["dataset_id"] = dataset_id

        rag_config = {
            "top_k": 4,
            "score_threshold": 0.0,
            "retrieval_mode": "keyword",
            "enable_reranker": False,
            "enable_multi_query": False,
            "enable_hyde": False,
            "enable_query_decomposition": False,
        }

        for case in fixture_cases:
            parser_backend = str(case.get("parser_backend") or "auto")
            resp = api.multipart(
                "POST",
                "/api/v1/documents/upload",
                fields={
                    "dataset_id": dataset_id,
                    "parser_backend": parser_backend,
                    "chunk_strategy": "langchain_recursive",
                    "governance_enabled": "true",
                    "chunk_vector_enabled": "true",
                    "bm25_index_enabled": "true",
                    "kg_enabled": "false",
                    "event_vector_enabled": "false",
                    "entity_vector_enabled": "false",
                },
                file_path=Path(case["path"]),
            )
            record_step(steps, f"upload:{case['name']}", resp)
            ensure_success(f"upload:{case['name']}", resp)
            document_id = str((resp.body or {}).get("id") or (resp.body or {}).get("document_id") or "")
            if not document_id:
                raise RuntimeError(f"upload:{case['name']} missing document id")

            detail = wait_for_document_completed(api, steps=steps, filename=str(case["name"]), document_id=document_id, poll_timeout=args.poll_timeout)
            chunks_resp = api.json("GET", f"/api/v1/documents/{document_id}/chunks?limit=200")
            record_step(steps, f"chunks:{case['name']}", chunks_resp)
            ensure_success(f"chunks:{case['name']}", chunks_resp)
            parsed_resp = api.json("GET", f"/api/v1/documents/{document_id}/parsed-content?max_chars=12000")
            record_step(steps, f"parsed:{case['name']}", parsed_resp)
            ensure_success(f"parsed:{case['name']}", parsed_resp)

            retrieve_resp = api.json(
                "POST",
                "/api/v1/rag/retrieve-preview",
                payload={"query": case["query"], "dataset_id": dataset_id, "rag_config": rag_config},
            )
            record_step(steps, f"retrieve:{case['name']}", retrieve_resp)
            ensure_success(f"retrieve:{case['name']}", retrieve_resp)

            chat_resp = api.json(
                "POST",
                "/api/v1/chat",
                payload={
                    "message": case["query"],
                    "dataset_id": dataset_id,
                    "stream": False,
                    "rag_config": {**rag_config, "answer_mode": "extractive", "max_tokens": 300},
                },
            )
            record_step(steps, f"chat:{case['name']}", chat_resp)
            ensure_success(f"chat:{case['name']}", chat_resp)

            merged_text = "\n".join(
                filter(
                    None,
                    [
                        parsed_text_from_response(parsed_resp.body),
                        response_text_from_body(retrieve_resp.body),
                        response_text_from_body(chat_resp.body),
                    ],
                )
            )
            failures = evaluate_format_case(
                {
                    "name": case["name"],
                    "expected_document_id": document_id,
                    "expected_terms": case["expected_terms"],
                    "min_chunks": 1,
                    "min_parsed_chars": 20,
                    "min_citations": 1,
                },
                document_id=document_id,
                chunk_count=len((chunks_resp.body or {}).get("items") or (chunks_resp.body or {}).get("chunks") or []) if isinstance(chunks_resp.body, dict) else 0,
                parsed_chars=len(parsed_text_from_response(parsed_resp.body)),
                citation_doc_ids=citation_document_ids(retrieve_resp.body),
                citation_count=len((retrieve_resp.body or {}).get("citations") or []) if isinstance(retrieve_resp.body, dict) else 0,
                response_text=merged_text,
            )
            result = {
                "name": case["name"],
                "file_path": str(case["path"]),
                "document_id": document_id,
                "status": str(detail.get("status") or "").lower(),
                "chunk_count": len((chunks_resp.body or {}).get("items") or (chunks_resp.body or {}).get("chunks") or []) if isinstance(chunks_resp.body, dict) else 0,
                "parsed_chars": len(parsed_text_from_response(parsed_resp.body)),
                "retrieve_status": retrieve_resp.status,
                "retrieve_citation_document_ids": citation_document_ids(retrieve_resp.body),
                "chat_status": chat_resp.status,
                "ok": not failures,
                "failures": failures,
            }
            summary["cases"].append(result)
            if failures:
                raise RuntimeError(f"format case failed {case['name']}: {failures}")

        summary["cleanup"] = cleanup_dataset(api, steps=steps, dataset_id=dataset_id)
        summary["ok"] = True
        return_code = 0
    except Exception as exc:  # noqa: BLE001
        summary["ok"] = False
        summary["error"] = str(exc)
        return_code = 1
    finally:
        report = {"summary": summary, "steps": steps}
        (artifact_dir / "report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return return_code


if __name__ == "__main__":
    raise SystemExit(main())
