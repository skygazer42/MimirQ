"""
PPTX parser (lightweight).

Extracts slide text (and simple tables) using python-pptx.

This is used as a robust fallback when general converters (MarkItDown/Pandoc)
fail on certain presentations.
"""


from pathlib import Path

from langchain_core.documents import Document
from pptx import Presentation

from app.rag.core.logging import get_logger

logger = get_logger(__name__)


def _clean_line(text: str) -> str:
    text = (text or "").replace("\r", " ").strip()
    return " ".join(text.split())


def _md_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [r + [""] * (width - len(r)) for r in rows]
    header = norm[0]
    body = norm[1:]
    out: list[str] = []
    out.append("| " + " | ".join(header) + " |")
    out.append("| " + " | ".join(["---"] * width) + " |")
    for row in body:
        out.append("| " + " | ".join(row) + " |")
    return "\n".join(out).strip()


def _escape_md_cell(text: str) -> str:
    return _clean_line(text).replace("|", r"\|")


def _extract_table_text(shape: object) -> str | None:
    try:
        if not bool(getattr(shape, "has_table", False)):
            return None
        table = shape.table
        rows: list[list[str]] = []
        for row_index in range(len(table.rows)):
            row: list[str] = []
            for col_index in range(len(table.columns)):
                try:
                    cell_text = table.cell(row_index, col_index).text
                except Exception:
                    cell_text = ""
                row.append(_escape_md_cell(cell_text))
            if any(row):
                rows.append(row)
        return _md_table(rows) if rows else ""
    except Exception as exc:
        # Best-effort: ignore table extraction errors.
        logger.debug("Failed to extract PPTX table; falling back to text frame handling: %s", exc)
        return None


def _paragraph_text(para: object) -> str:
    raw = str(getattr(para, "text", "") or "")
    text = _clean_line(raw)
    if not text:
        return ""
    try:
        level = int(getattr(para, "level", 0) or 0)
    except Exception:
        level = 0
    level = max(0, min(level, 10))
    prefix = ("  " * level) + "- " if text else ""
    return f"{prefix}{text}" if prefix else text


def _extract_shape_text(shape: object) -> list[str]:
    try:
        if not bool(getattr(shape, "has_text_frame", False)):
            return []
        tf = shape.text_frame
        parts: list[str] = []
        for para in getattr(tf, "paragraphs", []) or []:
            text = _paragraph_text(para)
            if text:
                parts.append(text)
        return parts
    except Exception:
        get_logger(__name__).debug("Skipping item after non-critical exception", exc_info=True)
        return []


def _extract_slide_parts(slide: object) -> list[str]:
    parts: list[str] = []
    for shape in getattr(slide, "shapes", []) or []:
        table_text = _extract_table_text(shape)
        if table_text is not None:
            if table_text:
                parts.append(table_text)
            continue
        parts.extend(_extract_shape_text(shape))
    return parts


class PptxParser:
    """Parse a PowerPoint .pptx into slide-level Documents."""

    def parse(self, file_path: Path) -> list[Document]:
        ext = file_path.suffix.lower()
        if ext != ".pptx":
            raise ValueError(f"PptxParser supports only .pptx, got: {ext or '(no ext)'}")

        prs = Presentation(str(file_path))
        total_slides = len(getattr(prs, "slides", []) or [])

        documents: list[Document] = []

        for idx, slide in enumerate(prs.slides):
            parts = _extract_slide_parts(slide)
            content = "\n".join([p for p in (parts or []) if p]).strip()
            if not content:
                continue

            documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": str(file_path.name),
                        "page": idx + 1,
                        "total_pages": total_slides,
                        "file_type": "pptx",
                    },
                )
            )

        if documents:
            return documents

        return [
            Document(
                page_content="",
                metadata={
                    "source": str(file_path.name),
                    "total_pages": total_slides,
                    "file_type": "pptx",
                },
            )
        ]
