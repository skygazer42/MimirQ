from pathlib import Path

import pytest


def _write_minimal_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello"
    slide.placeholders[1].text = "This is a test slide."
    prs.save(path)


def test_pptx_backend_parses_slide_text(tmp_path: Path) -> None:
    from app.parsing.factory import ParserFactory

    pptx_path = tmp_path / "sample.pptx"
    _write_minimal_pptx(pptx_path)

    factory = ParserFactory()
    docs, backend = factory.parse(pptx_path, parser_backend="pptx")
    assert backend == "pptx"
    assert docs
    assert any("Hello" in (d.page_content or "") for d in docs)


def test_markitdown_failure_falls_back_to_pptx(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    from app.parsing.factory import ParserFactory

    pptx_path = tmp_path / "sample.pptx"
    _write_minimal_pptx(pptx_path)

    factory = ParserFactory()

    class _Broken:
        def parse(self, *_args, **_kwargs):  # type: ignore[no-untyped-def]
            raise RuntimeError("boom")

    monkeypatch.setattr(factory, "_get_markitdown_parser", lambda: _Broken())

    docs, backend = factory.parse(pptx_path, parser_backend="auto")
    assert backend == "pptx"
    assert docs
    assert any("Hello" in (d.page_content or "") for d in docs)

